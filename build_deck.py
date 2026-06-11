# -*- coding: utf-8 -*-
"""Final presentation deck: Troubleshooting (병목 -> 해결) + Theory projection."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette (Ewha forest) ----
FOREST   = RGBColor(0x2C, 0x5F, 0x2D)
DARKGRN  = RGBColor(0x18, 0x33, 0x1A)
MOSS     = RGBColor(0x97, 0xBC, 0x62)
LIGHTMOSS= RGBColor(0xE8, 0xF0, 0xDC)
CREAM    = RGBColor(0xF6, 0xF7, 0xF2)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
INK      = RGBColor(0x2B, 0x2E, 0x2B)
MUTE     = RGBColor(0x6B, 0x72, 0x66)
TERRA    = RGBColor(0xB8, 0x50, 0x42)   # bottleneck red
TERRABG  = RGBColor(0xF4, 0xE6, 0xE3)
CARDLINE = RGBColor(0xDD, 0xE2, 0xD6)

FONT = "Malgun Gothic"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def set_ea(run, face=FONT):
    """Force East-Asian typeface so Korean renders in Malgun Gothic."""
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', face)


def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, wrap=True, space_after=None):
    """runs: list of paragraphs; each paragraph = list of (text, size, color, bold, italic)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space_after is not None:
            p.space_after = Pt(space_after)
        for (t, sz, col, b, it) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col
            r.font.bold = b; r.font.italic = it; r.font.name = FONT
            set_ea(r)
    return tb


def rect(slide, x, y, w, h, fill, line=None, line_w=None, rounded=False, shadow=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if rounded:
        try:
            shp.adjustments[0] = 0.06
        except Exception:
            pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    # shadow XML corrupts PowerPoint; use a faint border for card definition instead
    if line is not None:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w or 1)
    elif shadow:
        shp.line.color.rgb = CARDLINE; shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def oval(slide, x, y, d, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line: shp.line.color.rgb = line; shp.line.width = Pt(1.5)
    else: shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def arrow(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp


def kicker(slide, text, color=FOREST):
    add_text(slide, Inches(0.7), Inches(0.42), Inches(8), Inches(0.3),
             [[(text, 12, color, True, False)]])


def title(slide, text, color=INK):
    add_text(slide, Inches(0.7), Inches(0.72), Inches(12), Inches(0.8),
             [[(text, 30, color, True, False)]])


# =========================================================
# SLIDE 1 — Title (dark)
# =========================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, DARKGRN)
# moss side band
rect(s, 0, 0, Inches(0.28), SH, MOSS)
add_text(s, Inches(0.9), Inches(1.0), Inches(11), Inches(0.4),
         [[("이화여자대학교  ·  Data Engineering  ·  04 Methodology", 13, MOSS, True, False)]])
add_text(s, Inches(0.9), Inches(2.25), Inches(11.6), Inches(2.2),
         [[("문헌 기반 인과 추론,", 40, WHITE, True, False)],
          [("세 가지 병목과 해결", 40, WHITE, True, False)]],
         space_after=6)
add_text(s, Inches(0.92), Inches(4.35), Inches(11), Inches(0.5),
         [[("RAG Troubleshooting  —  검색 · 청킹 · 변수 차원의 문제를 진단하고 재설계하다", 16, MOSS, False, True)]])
# bottom pipeline chips
chips = ["Semantic Scholar", "Overlap Chunking", "Theory Mapping (ECD)", "Two-axis Query", "RAGAS"]
cx = Inches(0.9)
for c in chips:
    w = Inches(0.18 + 0.105 * len(c))
    rect(s, cx, Inches(5.55), w, Inches(0.46), RGBColor(0x24,0x4A,0x26), rounded=True)
    add_text(s, cx, Inches(5.55), w, Inches(0.46), [[(c, 11, LIGHTMOSS, False, False)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cx = Emu(cx + w + Inches(0.18))

# =========================================================
# SLIDE 2 — Three bottlenecks
# =========================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, CREAM)
kicker(s, "TROUBLESHOOTING")
title(s, "기존 파이프라인에서 발견된 세 가지 병목")
add_text(s, Inches(0.7), Inches(1.32), Inches(12), Inches(0.4),
         [[("증상: ", 13, TERRA, True, False),
           ("RAG 문서 연관성 스코어가 지나치게 낮아 인과관계 추론이 작동하지 않음", 13, MUTE, False, False)]])

cards = [
    ("1", "논문 검색 (Search API)",
     "도메인 키워드(gamification, student learning)로 검색 시 개론적 리뷰 페이퍼만 수집 → 구체적 변수 관계와 매칭 실패"),
    ("2", "청킹 (Chunking)",
     "GROBID 기반 정밀 청킹 시도 → API 토큰 한도·처리 속도 문제로 운영상 비현실적"),
    ("3", "검색·임베딩 (RAG)",
     "행동변수 간 직접 관계(풀스크린 ↔ 퀴즈 정답률)를 다룬 문헌이 애초에 부재 → 쿼리 확장만으로 해결 불가"),
]
cw = Inches(3.95); gap = Inches(0.27); x0 = Inches(0.7); y0 = Inches(2.0); ch = Inches(4.05)
for i, (num, head, body) in enumerate(cards):
    x = Emu(x0 + i * (cw + gap))
    rect(s, x, y0, cw, ch, WHITE, rounded=True, shadow=True)
    rect(s, x, y0, cw, Inches(0.14), MOSS if i < 2 else TERRA)
    oval(s, Emu(x + Inches(0.35)), Inches(2.45), Inches(0.85), LIGHTMOSS if i < 2 else TERRABG)
    add_text(s, Emu(x + Inches(0.35)), Inches(2.45), Inches(0.85), Inches(0.85),
             [[(num, 30, FOREST if i < 2 else TERRA, True, False)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Emu(x + Inches(0.35)), Inches(3.55), Emu(cw - Inches(0.7)), Inches(0.8),
             [[(head, 17, INK, True, False)]])
    add_text(s, Emu(x + Inches(0.35)), Inches(4.35), Emu(cw - Inches(0.7)), Inches(1.5),
             [[(body, 13, MUTE, False, False)]])

# =========================================================
# SLIDE 3 — Diagnosis: variable-dimension problem
# =========================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, WHITE)
kicker(s, "ROOT CAUSE")
title(s, "쿼리의 문제가 아니라 “변수 차원”의 문제")

# left insight panel
rect(s, Inches(0.7), Inches(1.7), Inches(5.75), Inches(5.1), CREAM, rounded=True)
add_text(s, Inches(1.05), Inches(2.0), Inches(5.1), Inches(0.5),
         [[("왜 문헌이 안 잡히는가", 18, FOREST, True, False)]])
pts = [
    ("행동변수 A–B의 직접 인과", "  (예: 풀스크린 FS → 퀴즈 정답률 PERF)를 검증한 논문은 사람이 찾아도 없다."),
    ("교육심리 SEM의 단위", "  심리 척도의 잠재 구성개념(construct)이지, 구체적 행동 로그 변수가 아니다."),
    ("따라서 필요한 것", "  행동 로그 변수를 이론적 차원(latent construct)에 투영한 뒤 그 차원으로 검색한다."),
]
yy = Inches(2.65)
for head, body in pts:
    oval(s, Inches(1.05), Emu(yy + Inches(0.07)), Inches(0.16), MOSS)
    add_text(s, Inches(1.4), yy, Inches(4.8), Inches(1.3),
             [[(head, 14, INK, True, False), (body, 14, MUTE, False, False)]])
    yy = Emu(yy + Inches(1.32))

# right: two reference cards
add_text(s, Inches(6.85), Inches(1.7), Inches(5.8), Inches(0.4),
         [[("선행 연구도 같은 전략: 비정형/미관측을 잠재요인으로", 14, FOREST, True, False)]])
refs = [
    ("COAT", "NeurIPS 2024 · Discovery of the Hidden World with LLMs",
     "LLM이 비정형 데이터(리뷰 등)에서 잠재요인(맛·가격)을 추출·점수화 → 정형 변수로 변환 후 인과 발견."),
    ("ARCADIA", "Scalable Causal Discovery for Corporate",
     "LLM 추론으로 잠재 경제 변수를 proxy confounder로 도입 → 미관측 교란 보정, back-door 식별성 확보."),
]
ry = Inches(2.25)
for name, cite, body in refs:
    rect(s, Inches(6.85), ry, Inches(5.78), Inches(2.05), WHITE, line=CARDLINE, line_w=1, rounded=True, shadow=True)
    rect(s, Inches(6.85), ry, Inches(0.13), Inches(2.05), FOREST)
    add_text(s, Inches(7.2), Emu(ry + Inches(0.25)), Inches(5.2), Inches(0.5),
             [[(name, 19, FOREST, True, False), ("   " + cite, 11, MUTE, False, True)]])
    add_text(s, Inches(7.2), Emu(ry + Inches(0.85)), Inches(5.2), Inches(1.1),
             [[(body, 13.5, INK, False, False)]])
    ry = Emu(ry + Inches(2.25))

# =========================================================
# SLIDE 4 — Theory projection diagram (the requested one)
# =========================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, CREAM)
kicker(s, "KEY IDEA")
title(s, "행동 로그 변수를 이론적 차원에 투영한다")
add_text(s, Inches(0.7), Inches(1.34), Inches(12), Inches(0.4),
         [[("Stealth Assessment의 ", 13, MUTE, False, False),
           ("ECD(Evidence-Centered Design)", 13, FOREST, True, False),
           (" 프레임워크 차용 — DAG를 그리기 전에 변수별 이론 메타를 부착", 13, MUTE, False, False)]])

# column headers
add_text(s, Inches(1.2), Inches(2.15), Inches(3.3), Inches(0.4),
         [[("행동 로그 변수", 14, INK, True, False)]], align=PP_ALIGN.CENTER)
add_text(s, Inches(5.55), Inches(2.15), Inches(2.0), Inches(0.4),
         [[("Theory Mapping", 13, FOREST, True, False)]], align=PP_ALIGN.CENTER)
add_text(s, Inches(8.35), Inches(2.15), Inches(4.0), Inches(0.4),
         [[("잠재 구성개념 (construct)", 14, INK, True, False)]], align=PP_ALIGN.CENTER)

rows = [
    ("FS", "Fullscreen 사용", "Immersion", "EGameFlow", "behavioral_proxy"),
    ("PERF", "퀴즈 정답률", "Knowledge Improvement", "EGameFlow", "outcome"),
    ("DIV", "행동 다양성", "Autonomy Satisfaction", "BANGS / SDT", "behavioral_proxy"),
]
ry = Inches(2.7); rh = Inches(1.05); rgap = Inches(0.28)
for code, desc, con, theo, role in rows:
    # left var box
    rect(s, Inches(1.2), ry, Inches(3.3), rh, WHITE, line=CARDLINE, line_w=1, rounded=True, shadow=True)
    add_text(s, Inches(1.45), Emu(ry + Inches(0.16)), Inches(2.9), Inches(0.45),
             [[(code, 18, TERRA, True, False), ("   " + desc, 12.5, MUTE, False, False)]])
    add_text(s, Inches(1.45), Emu(ry + Inches(0.62)), Inches(2.9), Inches(0.35),
             [[("behavioral log", 10.5, MUTE, False, True)]])
    # arrow
    arrow(s, Inches(4.75), Emu(ry + Inches(0.33)), Inches(1.55), Inches(0.4), MOSS)
    add_text(s, Inches(4.75), Emu(ry - Inches(0.02)), Inches(1.55), Inches(0.3),
             [[(role, 9.5, FOREST, True, False)]], align=PP_ALIGN.CENTER)
    # right construct box
    rect(s, Inches(8.35), ry, Inches(3.75), rh, FOREST, rounded=True, shadow=True)
    add_text(s, Inches(8.6), Emu(ry + Inches(0.16)), Inches(3.3), Inches(0.45),
             [[(con, 16, WHITE, True, False)]])
    add_text(s, Inches(8.6), Emu(ry + Inches(0.62)), Inches(3.3), Inches(0.35),
             [[(theo, 11, LIGHTMOSS, False, False)]])
    ry = Emu(ry + rh + rgap)

# footnote band
rect(s, Inches(1.2), Inches(6.85), Inches(10.9), Inches(0.0001), CREAM)
add_text(s, Inches(1.2), Inches(6.78), Inches(11), Inches(0.5),
         [[("연구자가 준 기반 이론 + 척도로 LLM이 직접 매핑 → 부착된 잠재 메타는 이후 인과 추론·쿼리 생성 시 참조 (직접적 행동변수 논문이 없으므로)",
            11.5, MUTE, False, True)]])

# =========================================================
# SLIDE 5 — Solutions
# =========================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, WHITE)
kicker(s, "SOLUTIONS")
title(s, "세 병목에 대응하는 재설계")

sols = [
    ("01", "검색 API 교체", "OpenAlex → Semantic Scholar",
     "긴 확장 쿼리에 불리한 키워드(BM25) 매칭 대신, 학습형 relevance 랭킹·SPECTER 임베딩으로 의미 매칭을 상대적으로 더 지원."),
    ("02", "청킹 단순화", "GROBID → Overlap Window",
     "정밀 파싱 없이 오버랩 슬라이딩 윈도우만으로도 문서 질·연관성 스코어 유지(RAGAS로 확인) → 토큰·속도 제약 해소."),
    ("03", "쿼리 확장 (두 축)", "이론 검증 논문 + Theory 매핑",
     "축1: 이론의 가설경로·SEM 검증 논문을 쿼리 그룹으로 검색 후 중복 제거. 축2: 매핑된 construct로 RAG 쿼리, 약한 변수는 변수명+desc 동시 투입."),
]
cw = Inches(3.95); gap = Inches(0.27); x0 = Inches(0.7); y0 = Inches(1.95); ch = Inches(3.45)
for i, (num, head, sub, body) in enumerate(sols):
    x = Emu(x0 + i * (cw + gap))
    rect(s, x, y0, cw, ch, CREAM, rounded=True, shadow=True)
    add_text(s, Emu(x + Inches(0.35)), Emu(y0 + Inches(0.3)), Inches(2), Inches(0.7),
             [[(num, 34, MOSS, True, False)]])
    add_text(s, Emu(x + Inches(0.35)), Emu(y0 + Inches(1.05)), Emu(cw - Inches(0.7)), Inches(0.5),
             [[(head, 18, FOREST, True, False)]])
    add_text(s, Emu(x + Inches(0.35)), Emu(y0 + Inches(1.5)), Emu(cw - Inches(0.7)), Inches(0.5),
             [[(sub, 12.5, TERRA, True, False)]])
    add_text(s, Emu(x + Inches(0.35)), Emu(y0 + Inches(2.0)), Emu(cw - Inches(0.7)), Inches(1.3),
             [[(body, 12.5, INK, False, False)]])

# code-ish callout for build_query
rect(s, Inches(0.7), Inches(5.7), Inches(11.93), Inches(1.05), DARKGRN, rounded=True)
add_text(s, Inches(1.0), Inches(5.85), Inches(11.4), Inches(0.4),
         [[("build_query()", 13, MOSS, True, False),
           ("   construct 단위로 인과 쿼리 생성", 11.5, LIGHTMOSS, False, True)]])
add_text(s, Inches(1.0), Inches(6.25), Inches(11.4), Inches(0.45),
         [[('f"{ca} & {cb} causal relationship in {learning_env} of {learner_type} learning"',
            12.5, WHITE, False, False)]])

# =========================================================
# SLIDE 6 — Evaluation design (RAGAS focused)
# =========================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, CREAM)
kicker(s, "EVALUATION")
title(s, "무엇을 평가할 것인가 — RAGAS, 선별적으로")
add_text(s, Inches(0.7), Inches(1.34), Inches(12), Inches(0.4),
         [[("유사도(연관성)만으론 부족 — 주제적으로 가까워도 인과 근거를 담았는지는 알 수 없다. 평가를 두 목적으로 분리한다.",
            13, MUTE, False, True)]])

# table
hdr = ["검증 목적", "지표", "채택", "이유"]
data = [
    ("청킹 교체가 검색 품질을 떨어뜨렸나", "Context Precision", "채택", "검색된 청크가 실제 근거를 담는지 — 청킹 ablation의 핵심"),
    ("(검색 품질, 회수 측면)", "Context Recall", "채택", "필요한 근거를 빠짐없이 회수했는지 (※ ground truth 필요)"),
    ("원래 병목 해소 확인", "Faithfulness", "채택", "인과 판단이 검색 텍스트에 실제 근거하는지 — 우리의 실패 모드"),
    ("QA 적합도", "Answer Relevancy", "제외", "과제가 QA가 아니라 ‘엣지 지지 여부’라 부적합"),
]
tx = Inches(0.7); ty = Inches(2.0); tw = Inches(11.93)
colw = [Inches(3.5), Inches(2.3), Inches(1.1), Inches(5.03)]
# header row
hx = tx
rect(s, tx, ty, tw, Inches(0.55), FOREST)
for j, htext in enumerate(hdr):
    add_text(s, Emu(hx + Inches(0.2)), ty, Emu(colw[j] - Inches(0.2)), Inches(0.55),
             [[(htext, 13, WHITE, True, False)]], anchor=MSO_ANCHOR.MIDDLE)
    hx = Emu(hx + colw[j])
# body rows
ry = Emu(ty + Inches(0.55)); rh = Inches(0.92)
for i, row in enumerate(data):
    bg = WHITE if i % 2 == 0 else LIGHTMOSS
    rect(s, tx, ry, tw, rh, bg)
    hx = tx
    for j, cell in enumerate(row):
        if j == 2:  # 채택/제외 badge
            ok = (cell == "채택")
            bw = Inches(0.78)
            rect(s, Emu(hx + Inches(0.16)), Emu(ry + Inches(0.26)), bw, Inches(0.4),
                 MOSS if ok else TERRA, rounded=True)
            add_text(s, Emu(hx + Inches(0.16)), Emu(ry + Inches(0.26)), bw, Inches(0.4),
                     [[(cell, 11.5, WHITE if not ok else DARKGRN, True, False)]],
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        else:
            col = INK if j != 0 else MUTE
            bold = (j == 1)
            add_text(s, Emu(hx + Inches(0.2)), ry, Emu(colw[j] - Inches(0.35)), rh,
                     [[(cell, 12.5 if j != 0 else 12, col, bold, j == 0)]],
                     anchor=MSO_ANCHOR.MIDDLE)
        hx = Emu(hx + colw[j])
    ry = Emu(ry + rh)

# note
rect(s, Inches(0.7), Inches(6.55), Inches(11.93), Inches(0.6), TERRABG, rounded=True)
add_text(s, Inches(1.0), Inches(6.55), Inches(11.4), Inches(0.6),
         [[("실무 주의: ", 12, TERRA, True, False),
           ("Context Recall은 정답 레퍼런스가 필요 → 소규모 골드셋을 구축하거나 RAGAS의 reference-free(LLM-judge) 변형 사용을 발표에 명시.",
            12, INK, False, False)]],
         anchor=MSO_ANCHOR.MIDDLE)

out = "Troubleshooting_Methodology.pptx"
prs.save(out)
print("saved", out, "slides", len(prs.slides._sldIdLst))
