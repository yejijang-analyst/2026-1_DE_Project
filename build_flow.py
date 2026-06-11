# -*- coding: utf-8 -*-
"""Minimal B&W pipeline flowchart, ~1/4 of a standard 16:9 slide.
Output = native, editable PowerPoint shapes (text boxes, connectors, triangles)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT  = "Consolas"

prs = Presentation()                       # default 13.333 x 7.5 (16:9)
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
s = prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE

# diagram occupies roughly the top-left quarter of the slide
CX = 3.55                 # diagram center x
LX, RX = 2.05, 5.05       # left / right column centers
NW = 2.45                 # node text width


def node(text, cx, y, size=8.5, bold=False, w=NW):
    tb = s.shapes.add_textbox(Inches(cx - w / 2), Inches(y), Inches(w), Inches(0.26))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_top = 0; tf.margin_bottom = 0; tf.margin_left = 0; tf.margin_right = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = BLACK


def vline(x, y1, y2):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y1), Inches(x), Inches(y2))
    c.line.color.rgb = BLACK; c.line.width = Pt(1)


def hline(x1, x2, y):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y), Inches(x2), Inches(y))
    c.line.color.rgb = BLACK; c.line.width = Pt(1)


def arrow_down(cx, y_tip):
    ww, hh = 0.16, 0.13
    sh = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                            Inches(cx - ww / 2), Inches(y_tip - hh), Inches(ww), Inches(hh))
    sh.rotation = 180
    sh.fill.solid(); sh.fill.fore_color.rgb = BLACK
    sh.line.fill.background(); sh.shadow.inherit = False


# ---- top: Researcher Metadata ----
node("Researcher Metadata", CX, 1.05, size=10, bold=True, w=4.6)
node("(theories, scales, domains)", CX, 1.33, size=8, w=4.6)

# split
vline(CX, 1.62, 1.82)
hline(LX, RX, 1.82)
vline(LX, 1.82, 2.06); arrow_down(LX, 2.08)
vline(RX, 1.82, 2.06); arrow_down(RX, 2.08)

# ---- left pipeline ----
left = [
    ("Search Query Expansion (NEW)", 2.14, True),
    ("Semantic Scholar API",         2.54, False),
    ("Academic Literature",          2.94, False),
    ("Chunking",                     3.34, False),
    ("ChromaDB + RAG",               3.74, False),
    ("Evidence Retrieval",           4.14, False),
]
for txt, y, b in left:
    node(txt, LX, y, bold=b)
for i in range(len(left) - 1):
    vline(LX, left[i][1] + 0.27, left[i + 1][1] - 0.01)

# ---- right column ----
node("Theory Mapping by LLM (NEW)", RX, 2.14, bold=True)
node("(Zero-shot assignment)",      RX, 2.42, size=8)
node("Latent Constructs",           RX, 4.14)
vline(RX, 2.72, 4.13)

# ---- merge ----
vline(LX, 4.41, 4.62)
vline(RX, 4.41, 4.62)
hline(LX, RX, 4.62)
vline(CX, 4.62, 4.86); arrow_down(CX, 4.88)

node("Construct-level Annotation", CX, 4.92, size=9, bold=True, w=4.6)
vline(CX, 5.2, 5.42); arrow_down(CX, 5.44)
node("DAG Refinement", CX, 5.48, size=9, bold=True, w=4.6)

out = "Pipeline_Flow_BW.pptx"
prs.save(out)
print("saved", out)
